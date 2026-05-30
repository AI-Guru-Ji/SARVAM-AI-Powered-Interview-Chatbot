"""
generate_technical_ppt.py — One-shot generator for the technical PPT.

Produces ``output/ShramSaathi_AI_Technical_Deck.pptx`` — a fully editable
PowerPoint deck describing the architecture, AI pipeline, services,
state machine and reliability patterns of ShramSaathi AI.

Run:
    python tools/generate_technical_ppt.py
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

# Make project imports work when run as a script.
ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

from constants.app_constants import (  # noqa: E402
    APP_NAME,
    APP_TAGLINE,
    COLOR_ACCENT,
    COLOR_DANGER,
    COLOR_NEUTRAL_100,
    COLOR_NEUTRAL_200,
    COLOR_NEUTRAL_500,
    COLOR_NEUTRAL_700,
    COLOR_NEUTRAL_900,
    COLOR_PRIMARY,
    COLOR_PRIMARY_DARK,
    COLOR_SUCCESS,
    COLOR_WARN,
    LLM_MODEL,
    STT_MODEL,
    TTS_MODEL,
)


# ──────────────────────────────────────────────────────────────────────
# Helpers — colours, slide dimensions, builders
# ──────────────────────────────────────────────────────────────────────
def hex_rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


PRIMARY = hex_rgb(COLOR_PRIMARY)
PRIMARY_DARK = hex_rgb(COLOR_PRIMARY_DARK)
ACCENT = hex_rgb(COLOR_ACCENT)
SUCCESS = hex_rgb(COLOR_SUCCESS)
WARN = hex_rgb(COLOR_WARN)
DANGER = hex_rgb(COLOR_DANGER)
NEUTRAL_100 = hex_rgb(COLOR_NEUTRAL_100)
NEUTRAL_200 = hex_rgb(COLOR_NEUTRAL_200)
NEUTRAL_500 = hex_rgb(COLOR_NEUTRAL_500)
NEUTRAL_700 = hex_rgb(COLOR_NEUTRAL_700)
NEUTRAL_900 = hex_rgb(COLOR_NEUTRAL_900)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Widescreen 16:9, 13.33 × 7.5 inches.
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _add_blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # 6 = blank


def _set_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _set_outline(shape, color: RGBColor, weight_pt: float = 0.75) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(weight_pt)


def _add_text(
    slide,
    text: str,
    *,
    left, top, width, height,
    font_size: int = 18,
    bold: bool = False,
    color: RGBColor = NEUTRAL_900,
    font_name: str = "Inter",
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _add_bullets(
    slide,
    bullets: list[str],
    *,
    left, top, width, height,
    font_size: int = 16,
    color: RGBColor = NEUTRAL_700,
    bullet_char: str = "•",
    line_spacing: float = 1.25,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = f"{bullet_char}  {text}"
        run.font.name = "Inter"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return box


def _add_brand_strip(slide) -> None:
    """Thin indigo strip across the top + small wordmark."""
    strip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.18),
    )
    _set_fill(strip, PRIMARY)
    _add_text(
        slide, APP_NAME,
        left=Inches(0.5), top=Inches(0.32),
        width=Inches(6), height=Inches(0.4),
        font_size=14, bold=True, color=PRIMARY_DARK,
    )
    _add_text(
        slide, "Technical Overview",
        left=Inches(0.5), top=Inches(0.6),
        width=Inches(6), height=Inches(0.3),
        font_size=10, color=NEUTRAL_500,
    )


def _add_page_footer(slide, page_num: int, total: int) -> None:
    _add_text(
        slide, f"{page_num} / {total}",
        left=Inches(12.3), top=Inches(7.1),
        width=Inches(0.8), height=Inches(0.3),
        font_size=10, color=NEUTRAL_500, align=PP_ALIGN.RIGHT,
    )
    _add_text(
        slide, "Built on Sarvam AI · ShramSaathi AI",
        left=Inches(0.5), top=Inches(7.1),
        width=Inches(6), height=Inches(0.3),
        font_size=10, color=NEUTRAL_500,
    )


def _add_section_header(slide, eyebrow: str, title: str) -> None:
    _add_brand_strip(slide)
    _add_text(
        slide, eyebrow.upper(),
        left=Inches(0.6), top=Inches(0.95),
        width=Inches(12), height=Inches(0.35),
        font_size=12, bold=True, color=ACCENT,
    )
    _add_text(
        slide, title,
        left=Inches(0.6), top=Inches(1.3),
        width=Inches(12.2), height=Inches(0.8),
        font_size=32, bold=True, color=NEUTRAL_900,
    )
    # Accent rule under the title
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.6), Inches(2.05), Inches(0.9), Inches(0.06),
    )
    _set_fill(rule, PRIMARY)


def _add_card(
    slide,
    *,
    left, top, width, height,
    fill: RGBColor = WHITE,
    border: RGBColor = NEUTRAL_200,
):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
    )
    card.adjustments[0] = 0.06
    _set_fill(card, fill)
    _set_outline(card, border, 0.75)
    card.shadow.inherit = False
    return card


def _add_chip(
    slide, text: str,
    *,
    left, top, width, height,
    fill: RGBColor, text_color: RGBColor = WHITE,
):
    chip = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
    )
    chip.adjustments[0] = 0.5
    _set_fill(chip, fill)
    chip.line.fill.background()
    tf = chip.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Inter"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = text_color
    return chip


def _add_arrow(slide, left, top, width, height, color: RGBColor = PRIMARY):
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, left, top, width, height,
    )
    _set_fill(arrow, color)
    return arrow


# ──────────────────────────────────────────────────────────────────────
# Slide builders
# ──────────────────────────────────────────────────────────────────────
def build_title_slide(prs: Presentation, total: int) -> None:
    s = _add_blank_slide(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _set_fill(bg, PRIMARY_DARK)
    # Accent bar
    bar = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.7), Inches(0.16), Inches(2.2),
    )
    _set_fill(bar, ACCENT)
    _add_text(
        s, "TECHNICAL DECK",
        left=Inches(1.1), top=Inches(2.7),
        width=Inches(8), height=Inches(0.5),
        font_size=14, bold=True, color=ACCENT,
    )
    _add_text(
        s, APP_NAME,
        left=Inches(1.1), top=Inches(3.0),
        width=Inches(11), height=Inches(1.4),
        font_size=72, bold=True, color=WHITE,
    )
    _add_text(
        s, "Voice-first interview bot for India's blue-collar workforce.",
        left=Inches(1.1), top=Inches(4.55),
        width=Inches(11), height=Inches(0.5),
        font_size=22, color=WHITE,
    )
    _add_text(
        s, APP_TAGLINE,
        left=Inches(1.1), top=Inches(5.05),
        width=Inches(11), height=Inches(0.5),
        font_size=18, color=hex_rgb("#C7D2FE"),
    )
    # Tech badges
    badges = [
        ("Sarvam AI", PRIMARY), ("Streamlit", ACCENT),
        ("Pydantic v2", SUCCESS), ("Clean Arch", WARN),
        ("6 Languages", DANGER),
    ]
    x = Inches(1.1)
    for label, col in badges:
        _add_chip(s, label, left=x, top=Inches(6.2),
                  width=Inches(1.7), height=Inches(0.4), fill=col)
        x += Inches(1.85)
    _add_text(
        s, "Built on Sarvam AI · 6 Indian Languages · ATS-ready output",
        left=Inches(1.1), top=Inches(6.9),
        width=Inches(11), height=Inches(0.4),
        font_size=12, color=hex_rgb("#A5B4FC"),
    )


def build_problem_slide(prs: Presentation, page: int, total: int) -> None:
    s = _add_blank_slide(prs)
    _add_section_header(s, "01 · Problem", "Hiring's literacy and language gap")

    bullets = [
        "Most blue-collar candidates in India have limited literacy.",
        "ATS / HR systems expect typed English input.",
        "Phone screens don't scale — recruiters can't interview every applicant.",
        "Existing voice tools don't support Indian languages well.",
    ]
    _add_bullets(
        s, bullets,
        left=Inches(0.7), top=Inches(2.4),
        width=Inches(6.5), height=Inches(4),
        font_size=18,
    )

    # Right side: the gap visual
    _add_card(s,
              left=Inches(7.6), top=Inches(2.4),
              width=Inches(5.1), height=Inches(4),
              fill=NEUTRAL_100, border=NEUTRAL_200)
    _add_text(s, "The Gap",
              left=Inches(7.85), top=Inches(2.55),
              width=Inches(4.7), height=Inches(0.4),
              font_size=14, bold=True, color=PRIMARY_DARK)
    _add_text(s, "Candidate speaks",
              left=Inches(7.85), top=Inches(3.0),
              width=Inches(4.7), height=Inches(0.4),
              font_size=14, color=NEUTRAL_700)
    _add_text(s, "Hindi · Bengali · Telugu …",
              left=Inches(7.85), top=Inches(3.35),
              width=Inches(4.7), height=Inches(0.4),
              font_size=14, bold=True, color=NEUTRAL_900)
    _add_arrow(s, Inches(7.85), Inches(4.0), Inches(4.7), Inches(0.4), color=PRIMARY)
    _add_text(s, "Recruiter receives",
              left=Inches(7.85), top=Inches(4.6),
              width=Inches(4.7), height=Inches(0.4),
              font_size=14, color=NEUTRAL_700)
    _add_text(s, "Structured English JSON + PDF",
              left=Inches(7.85), top=Inches(4.95),
              width=Inches(4.7), height=Inches(0.4),
              font_size=14, bold=True, color=SUCCESS)

    _add_page_footer(s, page, total)


def build_solution_slide(prs: Presentation, page: int, total: int) -> None:
    s = _add_blank_slide(prs)
    _add_section_header(s, "02 · Solution", "End-to-end voice interview, three stages")

    stages = [
        ("Onboarding", "9 voice questions capture identity, experience, salary, availability.", PRIMARY),
        ("Resume", "ATS-ready English PDF auto-generated from the conversation.", ACCENT),
        ("Interview", "5 role-specific questions + dynamic follow-ups → scored report.", SUCCESS),
    ]
    x = Inches(0.6)
    w = Inches(4.0)
    for label, body, color in stages:
        _add_card(s, left=x, top=Inches(2.5), width=w, height=Inches(4))
        bar = s.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, Inches(2.5), w, Inches(0.18),
        )
        _set_fill(bar, color)
        _add_text(s, label,
                  left=x + Inches(0.3), top=Inches(2.85),
                  width=w - Inches(0.6), height=Inches(0.5),
                  font_size=22, bold=True, color=NEUTRAL_900)
        _add_text(s, body,
                  left=x + Inches(0.3), top=Inches(3.45),
                  width=w - Inches(0.6), height=Inches(2.5),
                  font_size=14, color=NEUTRAL_700)
        x += w + Inches(0.15)
    _add_page_footer(s, page, total)


def build_languages_roles_slide(prs: Presentation, page: int, total: int) -> None:
    s = _add_blank_slide(prs)
    _add_section_header(s, "03 · Coverage", "Six languages · Four roles")

    # Languages
    _add_text(s, "Languages",
              left=Inches(0.7), top=Inches(2.4),
              width=Inches(4), height=Inches(0.4),
              font_size=18, bold=True, color=NEUTRAL_900)
    langs = ["English", "Hindi", "Bengali", "Telugu", "Punjabi", "Gujarati"]
    x = Inches(0.7)
    y = Inches(3.0)
    for i, lang in enumerate(langs):
        _add_chip(s, lang, left=x, top=y,
                  width=Inches(1.7), height=Inches(0.45), fill=PRIMARY)
        x += Inches(1.85)
        if (i + 1) % 3 == 0:
            x = Inches(0.7); y += Inches(0.65)

    # Roles
    _add_text(s, "Roles",
              left=Inches(0.7), top=Inches(5.0),
              width=Inches(4), height=Inches(0.4),
              font_size=18, bold=True, color=NEUTRAL_900)
    roles = [("🧹", "Housekeeping"), ("💡", "Electrician"),
             ("🔧", "Plumber"), ("🛡", "Security Guard")]
    x = Inches(0.7)
    for emoji, role in roles:
        _add_card(s, left=x, top=Inches(5.5),
                  width=Inches(2.9), height=Inches(1.2))
        _add_text(s, emoji,
                  left=x + Inches(0.2), top=Inches(5.65),
                  width=Inches(0.7), height=Inches(0.9),
                  font_size=32, color=NEUTRAL_900)
        _add_text(s, role,
                  left=x + Inches(0.9), top=Inches(5.95),
                  width=Inches(1.9), height=Inches(0.5),
                  font_size=16, bold=True, color=NEUTRAL_900)
        x += Inches(3.05)

    _add_page_footer(s, page, total)


def build_architecture_slide(prs: Presentation, page: int, total: int) -> None:
    s = _add_blank_slide(prs)
    _add_section_header(s, "04 · Architecture",
                        "UI-agnostic core + swappable presentation layer")

    # Top: UI layer
    _add_card(s, left=Inches(0.7), top=Inches(2.4),
              width=Inches(11.9), height=Inches(0.9),
              fill=hex_rgb("#EEF2FF"), border=PRIMARY)
    _add_text(s, "ui/streamlit/   (replaceable: ui/fastapi/, ui/cli/, ui/react/…)",
              left=Inches(0.9), top=Inches(2.65),
              width=Inches(11.5), height=Inches(0.45),
              font_size=16, bold=True, color=PRIMARY_DARK)

    # Middle: services
    _add_card(s, left=Inches(0.7), top=Inches(3.45),
              width=Inches(11.9), height=Inches(1.6),
              fill=hex_rgb("#ECFEFF"), border=ACCENT)
    _add_text(s, "services/   (UI-agnostic business logic)",
              left=Inches(0.9), top=Inches(3.55),
              width=Inches(11.5), height=Inches(0.4),
              font_size=14, bold=True, color=hex_rgb("#0E7490"))
    pills = [
        "sarvam_stt", "sarvam_tts", "sarvam_llm",
        "profile_service", "resume_service", "evaluation_service",
        "decide_next_turn", "health_check",
    ]
    x = Inches(0.9); y = Inches(4.0)
    for i, name in enumerate(pills):
        _add_chip(s, name, left=x, top=y,
                  width=Inches(1.35), height=Inches(0.4), fill=ACCENT)
        x += Inches(1.43)
        if (i + 1) % 8 == 0:
            x = Inches(0.9); y += Inches(0.5)

    # Bottom: foundations
    _add_card(s, left=Inches(0.7), top=Inches(5.2),
              width=Inches(11.9), height=Inches(1.6),
              fill=NEUTRAL_100, border=NEUTRAL_200)
    _add_text(s, "models/  ·  prompts/  ·  data/  ·  constants/  ·  config/  ·  utils/",
              left=Inches(0.9), top=Inches(5.3),
              width=Inches(11.5), height=Inches(0.4),
              font_size=14, bold=True, color=NEUTRAL_700)
    _add_bullets(s, [
        "models/schemas.py — Pydantic v2 I/O types (Evaluation, Profile, Resume)",
        "prompts/*.py — every LLM prompt extracted from service code",
        "data/ — translated questions + messages (one dict per language)",
        "constants/app_constants.py — all URLs, model IDs, thresholds, colours",
    ], left=Inches(0.9), top=Inches(5.7),
        width=Inches(11.5), height=Inches(1.2), font_size=12)

    _add_text(s,
              "Rule: services/ never imports from ui/. "
              "A new UI slots in without touching business logic.",
              left=Inches(0.7), top=Inches(6.95),
              width=Inches(12), height=Inches(0.3),
              font_size=11, color=NEUTRAL_500)

    _add_page_footer(s, page, total)


def build_tech_stack_slide(prs: Presentation, page: int, total: int) -> None:
    s = _add_blank_slide(prs)
    _add_section_header(s, "05 · Tech Stack", "Built for production demo realism")

    rows = [
        ("Web UI", "Streamlit + audio_recorder_streamlit"),
        ("STT", f"Sarvam Saaras v3 ({STT_MODEL})"),
        ("TTS", f"Sarvam Bulbul v3 ({TTS_MODEL})"),
        ("LLM", f"Sarvam {LLM_MODEL} · reasoning_effort=low"),
        ("PDF", "fpdf2 (pure Python, no system deps)"),
        ("Config", "pydantic-settings (typed, .env-driven)"),
        ("Validation", "Pydantic v2 (typed I/O on every service)"),
        ("Tests", "pytest + unittest.mock (47 tests, ~0.4s)"),
    ]
    top = Inches(2.4)
    row_h = Inches(0.55)
    for i, (layer, tech) in enumerate(rows):
        y = top + row_h * i
        # Zebra stripe
        if i % 2 == 0:
            stripe = s.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.7), y, Inches(11.9), row_h,
            )
            _set_fill(stripe, NEUTRAL_100)
        _add_text(s, layer,
                  left=Inches(0.9), top=y + Inches(0.08),
                  width=Inches(3), height=Inches(0.4),
                  font_size=14, bold=True, color=PRIMARY_DARK)
        _add_text(s, tech,
                  left=Inches(4), top=y + Inches(0.08),
                  width=Inches(8.5), height=Inches(0.4),
                  font_size=14, color=NEUTRAL_700)

    _add_page_footer(s, page, total)


def build_ai_pipeline_slide(prs: Presentation, page: int, total: int) -> None:
    s = _add_blank_slide(prs)
    _add_section_header(s, "06 · AI Pipeline",
                        "Voice in → reasoning → voice out")

    # Three-stage pipeline visual
    stages = [
        ("🎙", "Saaras v3", "STT", "Indic audio → text in candidate language", ACCENT),
        ("🧠", "sarvam-30b", "LLM", "Reasoning · follow-ups · resume · evaluation", PRIMARY),
        ("🔊", "Bulbul v3", "TTS", "Question / closing audio in candidate language", SUCCESS),
    ]
    x = Inches(0.7)
    w = Inches(3.95)
    for i, (emoji, model, label, body, color) in enumerate(stages):
        _add_card(s, left=x, top=Inches(2.6), width=w, height=Inches(3.6))
        _add_text(s, emoji,
                  left=x + Inches(0.3), top=Inches(2.8),
                  width=Inches(1), height=Inches(0.8),
                  font_size=42)
        _add_text(s, label,
                  left=x + Inches(1.3), top=Inches(2.9),
                  width=Inches(2.4), height=Inches(0.4),
                  font_size=12, bold=True, color=color)
        _add_text(s, model,
                  left=x + Inches(1.3), top=Inches(3.2),
                  width=Inches(2.4), height=Inches(0.5),
                  font_size=20, bold=True, color=NEUTRAL_900)
        _add_text(s, body,
                  left=x + Inches(0.3), top=Inches(4.2),
                  width=w - Inches(0.6), height=Inches(2),
                  font_size=13, color=NEUTRAL_700)
        if i < 2:
            _add_arrow(s,
                       x + w + Inches(0.05), Inches(4.1),
                       Inches(0.5), Inches(0.4),
                       color=PRIMARY)
        x += w + Inches(0.6)

    # Resilience callouts
    _add_text(s, "Resilience built into every stage",
              left=Inches(0.7), top=Inches(6.4),
              width=Inches(12), height=Inches(0.35),
              font_size=14, bold=True, color=NEUTRAL_900)
    _add_text(
        s,
        "STT min-byte guard · LLM token-budget retry (Hindi → English) · "
        "partial-JSON regex recovery · deterministic fallbacks · pre-baked welcome audio",
        left=Inches(0.7), top=Inches(6.75),
        width=Inches(12), height=Inches(0.4),
        font_size=12, color=NEUTRAL_500,
    )

    _add_page_footer(s, page, total)


def build_fsm_slide(prs: Presentation, page: int, total: int) -> None:
    s = _add_blank_slide(prs)
    _add_section_header(s, "07 · State Machine",
                        "10-stage Streamlit FSM, dispatch by stage")

    stages = [
        ("setup", PRIMARY),
        ("profile_intro", PRIMARY),
        ("profile", ACCENT),
        ("profile_building", ACCENT),
        ("profile_review", ACCENT),
        ("greeting", WARN),
        ("interview", WARN),
        ("closing", WARN),
        ("evaluating", DANGER),
        ("report", SUCCESS),
    ]
    # Two-row layout
    row1 = stages[:5]
    row2 = stages[5:]
    chip_w = Inches(2.2)
    chip_h = Inches(0.5)
    gap = Inches(0.2)
    arrow_w = Inches(0.3)

    def render_row(row, y):
        x = Inches(0.7)
        for i, (name, color) in enumerate(row):
            _add_chip(s, name, left=x, top=y,
                      width=chip_w, height=chip_h, fill=color)
            x += chip_w
            if i < len(row) - 1:
                _add_arrow(s, x, y + Inches(0.1),
                           arrow_w, Inches(0.3), color=NEUTRAL_500)
                x += arrow_w + gap

    render_row(row1, Inches(2.7))
    render_row(row2, Inches(4.2))

    # Curving arrow note between rows
    _add_text(s, "↩  continues",
              left=Inches(10.5), top=Inches(3.4),
              width=Inches(2), height=Inches(0.4),
              font_size=12, color=NEUTRAL_500)

    bullets = [
        "Each stage owns one view file in ui/streamlit/views/.",
        "Stage transitions happen inside view callbacks → st.rerun().",
        "reset_interview() clears interview-level state and returns to setup.",
        "session-level state (welcome unlock, TTS cache) survives between candidates.",
    ]
    _add_bullets(s, bullets,
                 left=Inches(0.7), top=Inches(5.4),
                 width=Inches(12), height=Inches(2), font_size=14)

    _add_page_footer(s, page, total)


def build_services_slide(prs: Presentation, page: int, total: int) -> None:
    s = _add_blank_slide(prs)
    _add_section_header(s, "08 · Services",
                        "Where the work happens — small, typed, testable")

    services = [
        ("SarvamSttService", "Audio bytes → transcript", PRIMARY,
         "Saaras v3, language hinted by lang_code; min-byte audio guard"),
        ("SarvamTtsService", "Text → WAV bytes", PRIMARY,
         "Bulbul v3, speaker pre-selected, output cached at app level"),
        ("SarvamLlmService", "Chat/completions wrapper", PRIMARY,
         "Single retry-aware caller; surfaces LlmBudgetExceededError"),
        ("ProfileService", "9 answers → Profile JSON", ACCENT,
         "Deterministic Q→field mapping; zero LLM calls"),
        ("ResumeService", "Profile → ATS PDF", ACCENT,
         "LLM-driven; deterministic fallback when LLM fails"),
        ("EvaluationService", "Interview → scorecard", SUCCESS,
         "Retry in English on budget bust · partial-JSON regex recovery"),
        ("DecideNextTurnService", "Answer → next action", WARN,
         "Decides follow_up vs advance, max 1 follow-up per question"),
        ("HealthCheckService", "Pre-demo readiness", NEUTRAL_500,
         "Tiny probe to TTS + LLM with HEALTH_CHECK timeouts"),
    ]
    col_w = Inches(6.0)
    row_h = Inches(1.05)
    for i, (name, role, color, detail) in enumerate(services):
        col = i % 2
        row = i // 2
        x = Inches(0.6) + (col_w + Inches(0.2)) * col
        y = Inches(2.4) + row_h * row
        _add_card(s, left=x, top=y, width=col_w, height=row_h - Inches(0.1))
        # Color stripe on left
        stripe = s.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y, Inches(0.12), row_h - Inches(0.1),
        )
        _set_fill(stripe, color)
        _add_text(s, name,
                  left=x + Inches(0.25), top=y + Inches(0.1),
                  width=col_w - Inches(0.4), height=Inches(0.35),
                  font_size=14, bold=True, color=NEUTRAL_900)
        _add_text(s, role,
                  left=x + Inches(0.25), top=y + Inches(0.42),
                  width=col_w - Inches(0.4), height=Inches(0.3),
                  font_size=11, color=color)
        _add_text(s, detail,
                  left=x + Inches(0.25), top=y + Inches(0.66),
                  width=col_w - Inches(0.4), height=Inches(0.35),
                  font_size=11, color=NEUTRAL_500)

    _add_page_footer(s, page, total)


def build_reliability_slide(prs: Presentation, page: int, total: int) -> None:
    s = _add_blank_slide(prs)
    _add_section_header(s, "09 · Reliability",
                        "Battle-scars baked into the code")

    items = [
        ("⚠", "LLM token-budget retry",
         "Evaluation auto-retries in English (~3× more token-efficient) "
         "when Hindi/Devanagari busts the 4096-token cap."),
        ("🧯", "Partial-JSON recovery",
         "Regex extraction salvages scores + summary when the LLM "
         "truncates mid-JSON."),
        ("🧱", "Deterministic profile builder",
         "Zero LLM calls — the 9-question → field mapping is a "
         "fixed schema. No surprises."),
        ("⚡", "Pre-baked welcome audio",
         "assets/welcome.wav ships with the repo. Demo opens with "
         "zero Sarvam latency."),
        ("🔓", "Click-to-begin overlay",
         "Satisfies browser autoplay policy so the welcome voice "
         "always plays on the first interaction."),
        ("🗂", "Two-bucket session state",
         "reset_interview() wipes one candidate's data without "
         "clearing welcome unlock or TTS cache."),
        ("🧬", "Mutable-default deepcopy",
         "Every reset uses copy.deepcopy() so interview N+1 never "
         "sees interview N's transcripts."),
        ("🩹", "Graceful retry UI",
         "Failed evaluations show a 'Retry Evaluation' screen, "
         "never a Python traceback."),
    ]
    col_w = Inches(6.0)
    row_h = Inches(1.05)
    for i, (icon, title, body) in enumerate(items):
        col = i % 2
        row = i // 2
        x = Inches(0.6) + (col_w + Inches(0.2)) * col
        y = Inches(2.4) + row_h * row
        _add_card(s, left=x, top=y, width=col_w, height=row_h - Inches(0.1))
        _add_text(s, icon,
                  left=x + Inches(0.18), top=y + Inches(0.18),
                  width=Inches(0.6), height=Inches(0.6),
                  font_size=24)
        _add_text(s, title,
                  left=x + Inches(0.85), top=y + Inches(0.12),
                  width=col_w - Inches(1), height=Inches(0.35),
                  font_size=13, bold=True, color=NEUTRAL_900)
        _add_text(s, body,
                  left=x + Inches(0.85), top=y + Inches(0.42),
                  width=col_w - Inches(1), height=Inches(0.55),
                  font_size=10, color=NEUTRAL_500)

    _add_page_footer(s, page, total)


def build_state_model_slide(prs: Presentation, page: int, total: int) -> None:
    s = _add_blank_slide(prs)
    _add_section_header(s, "10 · State Model",
                        "Two-bucket session state, deep-copied on assignment")

    # SESSION level
    _add_card(s, left=Inches(0.6), top=Inches(2.4),
              width=Inches(5.9), height=Inches(4.3),
              fill=hex_rgb("#EEF2FF"), border=PRIMARY)
    _add_text(s, "SESSION_LEVEL  (preserved across interviews)",
              left=Inches(0.8), top=Inches(2.55),
              width=Inches(5.5), height=Inches(0.5),
              font_size=15, bold=True, color=PRIMARY_DARK)
    _add_bullets(s, [
        "welcome_unlocked — overlay shown once per tab",
        "setup_welcome_played — welcome plays once per tab",
        "_tts_cache — keyed by (text, language)",
    ], left=Inches(0.85), top=Inches(3.1),
        width=Inches(5.6), height=Inches(2),
        font_size=13, color=NEUTRAL_700)
    _add_text(s, "Why preserve?",
              left=Inches(0.8), top=Inches(5.1),
              width=Inches(5.5), height=Inches(0.3),
              font_size=12, bold=True, color=PRIMARY_DARK)
    _add_text(s,
              "Wiping these between interviews would replay the welcome "
              "on every 'New Interview' click and force a redundant Sarvam "
              "round-trip on every TTS call.",
              left=Inches(0.85), top=Inches(5.4),
              width=Inches(5.5), height=Inches(1.2),
              font_size=11, color=NEUTRAL_500)

    # INTERVIEW level
    _add_card(s, left=Inches(6.7), top=Inches(2.4),
              width=Inches(5.9), height=Inches(4.3),
              fill=hex_rgb("#FFF7ED"), border=WARN)
    _add_text(s, "INTERVIEW_LEVEL  (wiped on New / Abort)",
              left=Inches(6.9), top=Inches(2.55),
              width=Inches(5.5), height=Inches(0.5),
              font_size=15, bold=True, color=hex_rgb("#92400E"))
    _add_bullets(s, [
        "stage · role · lang · candidate_name",
        "question_idx · follow_up_count · transcripts",
        "profile_json · profile_resume · profile_resume_pdf",
        "evaluation · report_saved_path",
        "greeting / closing audio + play flags",
    ], left=Inches(6.95), top=Inches(3.1),
        width=Inches(5.6), height=Inches(2.6),
        font_size=13, color=NEUTRAL_700)
    _add_text(s, "Why deepcopy?",
              left=Inches(6.9), top=Inches(5.55),
              width=Inches(5.5), height=Inches(0.3),
              font_size=12, bold=True, color=hex_rgb("#92400E"))
    _add_text(s,
              "Mutable defaults like 'transcripts: []' would share one "
              "list object — interview N+1 would see N's history. "
              "copy.deepcopy() on every assignment fixes the aliasing.",
              left=Inches(6.95), top=Inches(5.85),
              width=Inches(5.5), height=Inches(1),
              font_size=11, color=NEUTRAL_500)

    _add_page_footer(s, page, total)


def build_evaluation_slide(prs: Presentation, page: int, total: int) -> None:
    s = _add_blank_slide(prs)
    _add_section_header(s, "11 · Evaluation",
                        "LLM scorecard with retry + partial-JSON recovery")

    # Flow: transcripts → attempt 1 → attempt 2 → recovery → report
    nodes = [
        ("Transcripts\n(Q + A pairs)", NEUTRAL_500),
        ("Attempt 1\n(native lang)", PRIMARY),
        ("Budget\nbusts?", WARN),
        ("Attempt 2\n(English)", ACCENT),
        ("Partial-JSON\nrecovery", DANGER),
        ("Scorecard\n+ summary", SUCCESS),
    ]
    x = Inches(0.5)
    w = Inches(1.85)
    for i, (label, color) in enumerate(nodes):
        _add_card(s, left=x, top=Inches(3.0),
                  width=w, height=Inches(1.4),
                  fill=color, border=color)
        tf_box = _add_text(s, label,
                           left=x, top=Inches(3.0),
                           width=w, height=Inches(1.4),
                           font_size=11, bold=True, color=WHITE,
                           align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < len(nodes) - 1:
            _add_arrow(s, x + w + Inches(0.02), Inches(3.5),
                       Inches(0.2), Inches(0.4), color=NEUTRAL_500)
        x += w + Inches(0.25)

    _add_text(s, "Score schema",
              left=Inches(0.7), top=Inches(5.0),
              width=Inches(12), height=Inches(0.4),
              font_size=16, bold=True, color=NEUTRAL_900)
    _add_bullets(s, [
        "overall_score (0–10) · technical_score · communication_score · attitude_score",
        "strengths · improvements · summary (one paragraph)",
        "generation_failed flag → UI shows Retry Evaluation screen, never a traceback",
    ], left=Inches(0.7), top=Inches(5.4),
        width=Inches(12), height=Inches(2),
        font_size=13)

    _add_page_footer(s, page, total)


def build_testing_slide(prs: Presentation, page: int, total: int) -> None:
    s = _add_blank_slide(prs)
    _add_section_header(s, "12 · Testing", "47 unit tests · ~0.4s · zero API calls")

    bullets = [
        "tests/conftest.py mocks every Sarvam endpoint via unittest.mock — no API key needed.",
        "Service tests cover happy paths plus the recovery branches "
        "(LLM budget bust, partial JSON, deterministic fallbacks).",
        "Schema tests assert Pydantic v2 round-trips and the field-order "
        "contract between QUESTION_FIELD_ORDER and PROFILE_QUESTIONS.",
        "Utility tests cover phone/email extraction in Hindi digit words, "
        "+91 prefixes, 'at' / 'dot' patterns.",
        "CI-friendly: pytest exits cleanly, no network, no secrets.",
    ]
    _add_bullets(s, bullets,
                 left=Inches(0.7), top=Inches(2.5),
                 width=Inches(7.5), height=Inches(4.5),
                 font_size=14)

    # Stat cards
    stats = [("47", "tests"), ("~0.4s", "runtime"), ("0", "API calls")]
    x = Inches(8.6)
    for value, label in stats:
        _add_card(s, left=x, top=Inches(2.5),
                  width=Inches(1.3), height=Inches(1.3))
        _add_text(s, value,
                  left=x, top=Inches(2.65),
                  width=Inches(1.3), height=Inches(0.6),
                  font_size=24, bold=True, color=PRIMARY,
                  align=PP_ALIGN.CENTER)
        _add_text(s, label,
                  left=x, top=Inches(3.25),
                  width=Inches(1.3), height=Inches(0.4),
                  font_size=11, color=NEUTRAL_500,
                  align=PP_ALIGN.CENTER)
        x += Inches(1.4)

    _add_text(s, "Commands",
              left=Inches(8.6), top=Inches(4.2),
              width=Inches(4), height=Inches(0.4),
              font_size=14, bold=True, color=NEUTRAL_900)
    cmd_box = _add_card(s, left=Inches(8.6), top=Inches(4.55),
                        width=Inches(4.2), height=Inches(1.8),
                        fill=NEUTRAL_900, border=NEUTRAL_900)
    cmd_tf = cmd_box.text_frame
    cmd_tf.word_wrap = True
    cmd_tf.margin_left = Inches(0.2)
    cmd_tf.margin_top = Inches(0.15)
    for cmd in ["$ pytest", "$ pytest -k recovery", "$ pytest -v"]:
        p = cmd_tf.add_paragraph() if cmd != "$ pytest" else cmd_tf.paragraphs[0]
        p.text = cmd
        for r in p.runs:
            r.font.name = "Consolas"
            r.font.size = Pt(13)
            r.font.color.rgb = hex_rgb("#10B981")

    _add_page_footer(s, page, total)


def build_extensibility_slide(prs: Presentation, page: int, total: int) -> None:
    s = _add_blank_slide(prs)
    _add_section_header(s, "13 · Extending the System",
                        "Where to change what")

    rows = [
        ("Add a role",
         "data/interview_questions.py → QUESTION_BANK (one entry × 6 languages)"),
        ("Add a language",
         "data/interview_questions.py — LANGUAGES + every translated dict"),
        ("Change an LLM prompt",
         "prompts/<feature>_prompt.py — never inline in services"),
        ("Tune a constant",
         "constants/app_constants.py (URLs, timeouts, model IDs, colours)"),
        ("Change welcome text",
         "Edit SETUP_WELCOME_TEXT, then python tools/regenerate_welcome.py"),
        ("Add an LLM service",
         "services/<feature>_service.py — accept Settings, raise from utils/exceptions"),
        ("Swap Streamlit for FastAPI / React",
         "Add ui/<your-stack>/. Call the same services. Done."),
    ]
    top = Inches(2.4)
    row_h = Inches(0.6)
    for i, (want, where) in enumerate(rows):
        y = top + row_h * i
        if i % 2 == 0:
            stripe = s.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.7), y, Inches(11.9), row_h,
            )
            _set_fill(stripe, NEUTRAL_100)
        _add_text(s, want,
                  left=Inches(0.9), top=y + Inches(0.12),
                  width=Inches(3.5), height=Inches(0.4),
                  font_size=14, bold=True, color=PRIMARY_DARK)
        _add_text(s, where,
                  left=Inches(4.6), top=y + Inches(0.12),
                  width=Inches(8), height=Inches(0.4),
                  font_size=13, color=NEUTRAL_700)

    _add_page_footer(s, page, total)


def build_config_run_slide(prs: Presentation, page: int, total: int) -> None:
    s = _add_blank_slide(prs)
    _add_section_header(s, "14 · Run It", "Setup → bake welcome → launch")

    steps = [
        ("1", "Install",
         "python -m venv venv && source venv/bin/activate\n"
         "pip install -r requirements.txt\n"
         "cp .env.example .env  # paste SARVAM_API_KEY"),
        ("2", "Bake welcome (one-time)",
         "python tools/regenerate_welcome.py"),
        ("3", "Launch",
         "streamlit run ui/streamlit/app.py\n"
         "# opens http://localhost:8501"),
        ("4", "Pre-demo health check (optional)",
         "python main.py --health-check"),
    ]
    top = Inches(2.4)
    for i, (num, title, body) in enumerate(steps):
        y = top + Inches(1.1) * i
        # Step number circle
        circle = s.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.7), y, Inches(0.8), Inches(0.8),
        )
        _set_fill(circle, PRIMARY)
        tf = circle.text_frame
        tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = num
        r.font.name = "Inter"; r.font.size = Pt(28); r.font.bold = True
        r.font.color.rgb = WHITE
        _add_text(s, title,
                  left=Inches(1.7), top=y,
                  width=Inches(11), height=Inches(0.4),
                  font_size=16, bold=True, color=NEUTRAL_900)
        # Code body
        body_box = s.shapes.add_textbox(
            Inches(1.7), y + Inches(0.4),
            Inches(11), Inches(0.7),
        )
        btf = body_box.text_frame
        btf.word_wrap = True
        for j, line in enumerate(body.split("\n")):
            p = btf.paragraphs[0] if j == 0 else btf.add_paragraph()
            r = p.add_run()
            r.text = line
            r.font.name = "Consolas"
            r.font.size = Pt(12)
            r.font.color.rgb = NEUTRAL_700

    _add_page_footer(s, page, total)


def build_roadmap_slide(prs: Presentation, page: int, total: int) -> None:
    s = _add_blank_slide(prs)
    _add_section_header(s, "15 · Roadmap", "What ships next")

    items = [
        ("NOW", "Working demo",
         "6 languages, 4 roles, end-to-end voice flow, ATS PDF, scored report.",
         SUCCESS),
        ("NEXT", "Persistence layer",
         "SQLite/Postgres for candidate history, recruiter dashboard, "
         "audit log of all answers.",
         WARN),
        ("LATER", "Multi-tenant deployment",
         "FastAPI service backing a React frontend; one schema per "
         "recruiter org; SSO; Sarvam quota management.",
         PRIMARY),
        ("FUTURE", "Skill-test integrations",
         "Optional plug-ins for role-specific skill tests "
         "(electrician circuit diagrams, plumber pipe layouts).",
         ACCENT),
    ]
    x = Inches(0.6)
    w = Inches(3.0)
    for label, title, body, color in items:
        _add_card(s, left=x, top=Inches(2.5), width=w, height=Inches(4))
        chip = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x + Inches(0.3), Inches(2.7),
            Inches(1.3), Inches(0.4),
        )
        chip.adjustments[0] = 0.5
        _set_fill(chip, color)
        chip.line.fill.background()
        ctf = chip.text_frame
        ctf.margin_top = Emu(0); ctf.margin_bottom = Emu(0)
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = label
        cr.font.name = "Inter"; cr.font.size = Pt(11); cr.font.bold = True
        cr.font.color.rgb = WHITE
        _add_text(s, title,
                  left=x + Inches(0.3), top=Inches(3.3),
                  width=w - Inches(0.6), height=Inches(0.6),
                  font_size=18, bold=True, color=NEUTRAL_900)
        _add_text(s, body,
                  left=x + Inches(0.3), top=Inches(4.0),
                  width=w - Inches(0.6), height=Inches(2.4),
                  font_size=12, color=NEUTRAL_700)
        x += w + Inches(0.15)

    _add_page_footer(s, page, total)


def build_thanks_slide(prs: Presentation, page: int, total: int) -> None:
    s = _add_blank_slide(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _set_fill(bg, PRIMARY_DARK)
    _add_text(s, "Thank you",
              left=Inches(1), top=Inches(2.8),
              width=Inches(11), height=Inches(1.4),
              font_size=72, bold=True, color=WHITE)
    _add_text(
        s, "ShramSaathi AI — voice-first hiring for India's blue-collar workforce.",
        left=Inches(1), top=Inches(4.3),
        width=Inches(11), height=Inches(0.5),
        font_size=20, color=hex_rgb("#C7D2FE"),
    )
    _add_text(
        s, "Built on Sarvam AI.",
        left=Inches(1), top=Inches(4.8),
        width=Inches(11), height=Inches(0.5),
        font_size=18, color=ACCENT,
    )


# ──────────────────────────────────────────────────────────────────────
# Main
# ──────────────────────────────────────────────────────────────────────
def main() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        build_title_slide,
        build_problem_slide,
        build_solution_slide,
        build_languages_roles_slide,
        build_architecture_slide,
        build_tech_stack_slide,
        build_ai_pipeline_slide,
        build_fsm_slide,
        build_services_slide,
        build_reliability_slide,
        build_state_model_slide,
        build_evaluation_slide,
        build_testing_slide,
        build_extensibility_slide,
        build_config_run_slide,
        build_roadmap_slide,
        build_thanks_slide,
    ]
    total = len(builders)

    # Title slide has its own signature; others take (prs, page, total).
    build_title_slide(prs, total)
    for i, fn in enumerate(builders[1:-1], start=2):
        fn(prs, i, total)
    build_thanks_slide(prs, total, total)

    out_dir = ROOT / "output"
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / "ShramSaathi_AI_Technical_Deck.pptx"
    prs.save(out_path)
    print(f"✅ Wrote {total} slides to: {out_path}")


if __name__ == "__main__":
    main()
