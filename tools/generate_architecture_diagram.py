"""
generate_architecture_diagram.py — One-shot generator for the
architecture diagram.

Produces:
  * output/ShramSaathi_AI_Architecture.pptx  (editable)
  * output/ShramSaathi_AI_Architecture.pdf   (shareable, via LibreOffice)
  * output/ShramSaathi_AI_Architecture.png   (embed-friendly, via LibreOffice)

Reuses the brand palette from constants/app_constants.py so the diagram
stays in sync with the app's design system.

Run:
    python tools/generate_architecture_diagram.py
"""

from __future__ import annotations

import shutil
import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

from constants.app_constants import (  # noqa: E402
    APP_NAME,
    COLOR_ACCENT,
    COLOR_DANGER,
    COLOR_NEUTRAL_100,
    COLOR_NEUTRAL_200,
    COLOR_NEUTRAL_500,
    COLOR_NEUTRAL_700,
    COLOR_NEUTRAL_900,
    COLOR_PRIMARY,
    COLOR_PRIMARY_DARK,
    COLOR_SUCCESS,
    COLOR_WARN,
)


def hex_rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


PRIMARY = hex_rgb(COLOR_PRIMARY)
PRIMARY_DARK = hex_rgb(COLOR_PRIMARY_DARK)
ACCENT = hex_rgb(COLOR_ACCENT)
SUCCESS = hex_rgb(COLOR_SUCCESS)
WARN = hex_rgb(COLOR_WARN)
DANGER = hex_rgb(COLOR_DANGER)
N100 = hex_rgb(COLOR_NEUTRAL_100)
N200 = hex_rgb(COLOR_NEUTRAL_200)
N500 = hex_rgb(COLOR_NEUTRAL_500)
N700 = hex_rgb(COLOR_NEUTRAL_700)
N900 = hex_rgb(COLOR_NEUTRAL_900)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def outline(shape, color: RGBColor, weight_pt: float = 0.75) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(weight_pt)


def no_outline(shape) -> None:
    shape.line.fill.background()


def text(
    slide, t: str, *,
    left, top, width, height,
    size: int = 12, bold: bool = False,
    color: RGBColor = N900,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    font_name: str = "Inter",
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = t
    r.font.name = font_name
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def card(slide, *, left, top, width, height,
         bg: RGBColor = WHITE, border: RGBColor = N200,
         radius: float = 0.06):
    c = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
    )
    c.adjustments[0] = radius
    fill(c, bg)
    outline(c, border, 0.75)
    return c


def pill(slide, t: str, *, left, top, width, height,
         bg: RGBColor, fg: RGBColor = WHITE, size: int = 10):
    p = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
    )
    p.adjustments[0] = 0.5
    fill(p, bg)
    no_outline(p)
    tf = p.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = t
    run.font.name = "Inter"
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = fg
    return p


def arrow_down(slide, x, y, length, color: RGBColor = N500):
    """Vertical down arrow centred at x, starting at y."""
    a = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW,
        x - Inches(0.12), y, Inches(0.24), length,
    )
    fill(a, color)
    no_outline(a)
    return a


def arrow_h(slide, x, y, length, color: RGBColor = N500, right: bool = True):
    shape = MSO_SHAPE.RIGHT_ARROW if right else MSO_SHAPE.LEFT_ARROW
    a = slide.shapes.add_shape(
        shape, x, y - Inches(0.12), length, Inches(0.24),
    )
    fill(a, color)
    no_outline(a)
    return a


def build_diagram() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Brand strip
    strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.18))
    fill(strip, PRIMARY); no_outline(strip)
    text(s, APP_NAME,
         left=Inches(0.5), top=Inches(0.28),
         width=Inches(8), height=Inches(0.35),
         size=14, bold=True, color=PRIMARY_DARK)
    text(s, "System Architecture — UI-agnostic core + swappable presentation layer",
         left=Inches(0.5), top=Inches(0.58),
         width=Inches(12), height=Inches(0.3),
         size=11, color=N500)

    # ── Layout grid ────────────────────────────────────────────────────
    # Three horizontal layers on the LEFT (~ 9.0in wide).
    # External Sarvam APIs on the RIGHT (~ 3.0in wide).
    LEFT = Inches(0.5)
    LEFT_W = Inches(9.0)
    RIGHT = Inches(9.85)
    RIGHT_W = Inches(3.0)

    # ── Layer 1: User + UI ────────────────────────────────────────────
    y1 = Inches(1.1)
    h1 = Inches(1.05)
    # User browser box
    card(s, left=LEFT, top=y1, width=Inches(2.0), height=h1,
         bg=N100, border=N200)
    text(s, "👤 Candidate",
         left=LEFT + Inches(0.15), top=y1 + Inches(0.12),
         width=Inches(1.7), height=Inches(0.35),
         size=12, bold=True, color=N900)
    text(s, "Browser · Mic · Speaker",
         left=LEFT + Inches(0.15), top=y1 + Inches(0.45),
         width=Inches(1.7), height=Inches(0.5),
         size=10, color=N500)

    # UI layer
    ui_left = LEFT + Inches(2.3)
    ui_w = Inches(6.7)
    card(s, left=ui_left, top=y1, width=ui_w, height=h1,
         bg=hex_rgb("#EEF2FF"), border=PRIMARY)
    text(s, "Presentation Layer  ·  ui/streamlit/",
         left=ui_left + Inches(0.2), top=y1 + Inches(0.1),
         width=ui_w - Inches(0.4), height=Inches(0.35),
         size=13, bold=True, color=PRIMARY_DARK)
    text(s, "Replaceable: ui/fastapi/  ·  ui/cli/  ·  ui/react/",
         left=ui_left + Inches(0.2), top=y1 + Inches(0.4),
         width=ui_w - Inches(0.4), height=Inches(0.25),
         size=9, color=N500)
    ui_chips = ["app.py (FSM)", "state.py", "styles.py",
                "components.py", "views/*"]
    cx = ui_left + Inches(0.2)
    cy = y1 + Inches(0.7)
    cw = Inches(1.22)
    for ch in ui_chips:
        pill(s, ch, left=cx, top=cy,
             width=cw, height=Inches(0.28), bg=PRIMARY, size=9)
        cx += cw + Inches(0.08)

    # User ↔ UI two-way arrow
    arrow_h(s, LEFT + Inches(2.0), y1 + Inches(0.52),
            Inches(0.3), color=N700, right=True)

    # ── Layer 2: Services ─────────────────────────────────────────────
    y2 = Inches(2.45)
    h2 = Inches(2.5)
    card(s, left=LEFT, top=y2, width=LEFT_W, height=h2,
         bg=hex_rgb("#ECFEFF"), border=ACCENT)
    text(s, "Services Layer  ·  services/   (UI-agnostic business logic)",
         left=LEFT + Inches(0.2), top=y2 + Inches(0.1),
         width=LEFT_W - Inches(0.4), height=Inches(0.35),
         size=13, bold=True, color=hex_rgb("#0E7490"))
    text(s, "Each service accepts Settings in __init__; raises typed exceptions from utils/exceptions.py",
         left=LEFT + Inches(0.2), top=y2 + Inches(0.4),
         width=LEFT_W - Inches(0.4), height=Inches(0.25),
         size=9, color=N500)

    # Two rows of service cards
    services_row1 = [
        ("SarvamSttService", "Saaras v3 · audio→text", PRIMARY),
        ("SarvamTtsService", "Bulbul v3 · text→WAV", PRIMARY),
        ("SarvamLlmService", "sarvam-30b chat", PRIMARY),
        ("ProfileService", "9 Q → Profile JSON", ACCENT),
    ]
    services_row2 = [
        ("ResumeService", "Profile → ATS PDF", ACCENT),
        ("EvaluationService", "Interview → scorecard", SUCCESS),
        ("DecideNextTurnService", "follow_up vs advance", WARN),
        ("HealthCheckService", "TTS + LLM probe", N500),
    ]
    inner_left = LEFT + Inches(0.2)
    inner_w = LEFT_W - Inches(0.4)
    col_gap = Inches(0.12)
    svc_w = (inner_w - col_gap * 3) / 4
    svc_h = Inches(0.78)

    def render_svc_row(items, y):
        x = inner_left
        for name, body, color in items:
            card(s, left=x, top=y, width=svc_w, height=svc_h,
                 bg=WHITE, border=N200)
            bar = s.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x, y, Inches(0.08), svc_h,
            )
            fill(bar, color); no_outline(bar)
            text(s, name,
                 left=x + Inches(0.18), top=y + Inches(0.08),
                 width=svc_w - Inches(0.25), height=Inches(0.32),
                 size=11, bold=True, color=N900)
            text(s, body,
                 left=x + Inches(0.18), top=y + Inches(0.42),
                 width=svc_w - Inches(0.25), height=Inches(0.32),
                 size=9, color=N500)
            x += svc_w + col_gap

    render_svc_row(services_row1, y2 + Inches(0.75))
    render_svc_row(services_row2, y2 + Inches(1.62))

    # UI ↔ Services arrow
    arrow_down(s, LEFT + Inches(5.65), Inches(2.2), Inches(0.25),
               color=N700)

    # ── Layer 3: Foundations ──────────────────────────────────────────
    y3 = Inches(5.1)
    h3 = Inches(1.65)
    card(s, left=LEFT, top=y3, width=LEFT_W, height=h3,
         bg=N100, border=N200)
    text(s, "Foundation Layer",
         left=LEFT + Inches(0.2), top=y3 + Inches(0.1),
         width=LEFT_W - Inches(0.4), height=Inches(0.35),
         size=13, bold=True, color=N700)

    foundations = [
        ("models/", "schemas.py — Pydantic v2 I/O", PRIMARY),
        ("prompts/", "every LLM prompt extracted", ACCENT),
        ("data/", "translated Q + messages", SUCCESS),
        ("constants/", "URLs, model IDs, colours", WARN),
        ("config/", "pydantic-settings · .env", DANGER),
        ("utils/", "logger, exceptions, helpers, pdf", N500),
    ]
    fx = LEFT + Inches(0.2)
    fw = (LEFT_W - Inches(0.4) - Inches(0.1) * 5) / 6
    fh = Inches(0.95)
    fy = y3 + Inches(0.5)
    for name, body, color in foundations:
        card(s, left=fx, top=fy, width=fw, height=fh,
             bg=WHITE, border=N200)
        bar = s.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, fx, fy, fw, Inches(0.08),
        )
        fill(bar, color); no_outline(bar)
        text(s, name,
             left=fx + Inches(0.1), top=fy + Inches(0.18),
             width=fw - Inches(0.2), height=Inches(0.3),
             size=11, bold=True, color=N900)
        text(s, body,
             left=fx + Inches(0.1), top=fy + Inches(0.5),
             width=fw - Inches(0.2), height=Inches(0.4),
             size=8, color=N500)
        fx += fw + Inches(0.1)

    # Services ↔ Foundations arrow
    arrow_down(s, LEFT + Inches(5.65), Inches(4.95), Inches(0.18),
               color=N700)

    # ── External: Sarvam APIs (right side) ────────────────────────────
    y_ext = Inches(1.1)
    h_ext = Inches(3.85)
    card(s, left=RIGHT, top=y_ext, width=RIGHT_W, height=h_ext,
         bg=hex_rgb("#FEF3C7"), border=WARN)
    text(s, "External · Sarvam AI",
         left=RIGHT + Inches(0.2), top=y_ext + Inches(0.1),
         width=RIGHT_W - Inches(0.4), height=Inches(0.35),
         size=13, bold=True, color=hex_rgb("#92400E"))
    text(s, "api.sarvam.ai",
         left=RIGHT + Inches(0.2), top=y_ext + Inches(0.42),
         width=RIGHT_W - Inches(0.4), height=Inches(0.25),
         size=9, color=N500, font_name="Consolas")

    endpoints = [
        ("/speech-to-text", "Saaras v3", PRIMARY),
        ("/text-to-speech", "Bulbul v3", PRIMARY),
        ("/v1/chat/completions", "sarvam-30b · reasoning", PRIMARY),
    ]
    ey = y_ext + Inches(0.85)
    for path, model, color in endpoints:
        card(s, left=RIGHT + Inches(0.2), top=ey,
             width=RIGHT_W - Inches(0.4), height=Inches(0.85),
             bg=WHITE, border=N200)
        bar = s.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, RIGHT + Inches(0.2), ey,
            Inches(0.08), Inches(0.85),
        )
        fill(bar, color); no_outline(bar)
        text(s, path,
             left=RIGHT + Inches(0.4), top=ey + Inches(0.12),
             width=RIGHT_W - Inches(0.6), height=Inches(0.3),
             size=10, bold=True, color=N900, font_name="Consolas")
        text(s, model,
             left=RIGHT + Inches(0.4), top=ey + Inches(0.44),
             width=RIGHT_W - Inches(0.6), height=Inches(0.3),
             size=9, color=N500)
        ey += Inches(0.95)

    # Services → External arrow
    arrow_h(s, LEFT + LEFT_W, y2 + Inches(1.25),
            Inches(0.35), color=WARN, right=True)
    text(s, "HTTPS",
         left=LEFT + LEFT_W + Inches(0.05), top=y2 + Inches(0.95),
         width=Inches(0.7), height=Inches(0.25),
         size=8, color=N500)

    # ── Output box (bottom-right) ─────────────────────────────────────
    y_out = Inches(5.1)
    card(s, left=RIGHT, top=y_out, width=RIGHT_W, height=Inches(1.65),
         bg=hex_rgb("#D1FAE5"), border=SUCCESS)
    text(s, "Outputs",
         left=RIGHT + Inches(0.2), top=y_out + Inches(0.1),
         width=RIGHT_W - Inches(0.4), height=Inches(0.35),
         size=13, bold=True, color=hex_rgb("#065F46"))
    text(s,
         "• profile_{name}_{ts}.json\n"
         "• resume_{name}_{ts}.pdf\n"
         "• report_{name}_{ts}.json",
         left=RIGHT + Inches(0.2), top=y_out + Inches(0.5),
         width=RIGHT_W - Inches(0.4), height=Inches(1.1),
         size=10, color=N700, font_name="Consolas")

    # Bottom footer note on the FSM
    text(s,
         "FSM:  setup → profile_intro → profile → profile_building → profile_review → "
         "greeting → interview → closing → evaluating → report",
         left=Inches(0.5), top=Inches(6.95),
         width=Inches(12.3), height=Inches(0.3),
         size=10, color=N500, font_name="Consolas")
    text(s,
         "Layer rule:  services/ never imports from ui/.  A new UI slots in without touching business logic.",
         left=Inches(0.5), top=Inches(7.2),
         width=Inches(12.3), height=Inches(0.25),
         size=9, color=N500)

    return prs


def export_pdf_and_png(pptx_path: Path, out_dir: Path) -> tuple[Path, Path | None]:
    soffice = shutil.which("libreoffice") or shutil.which("soffice")
    if not soffice:
        print("⚠ LibreOffice not found — skipping PDF/PNG export")
        return pptx_path, None

    # PDF
    subprocess.run(
        [soffice, "--headless", "--convert-to", "pdf",
         "--outdir", str(out_dir), str(pptx_path)],
        check=True, capture_output=True,
    )
    pdf_path = out_dir / (pptx_path.stem + ".pdf")

    # PNG (single slide → single PNG)
    subprocess.run(
        [soffice, "--headless", "--convert-to", "png",
         "--outdir", str(out_dir), str(pptx_path)],
        check=True, capture_output=True,
    )
    png_path = out_dir / (pptx_path.stem + ".png")
    return pdf_path, png_path


def main() -> None:
    prs = build_diagram()
    out_dir = ROOT / "output"
    out_dir.mkdir(parents=True, exist_ok=True)
    pptx_path = out_dir / "ShramSaathi_AI_Architecture.pptx"
    prs.save(pptx_path)
    print(f"✅ PPTX: {pptx_path}")
    pdf_path, png_path = export_pdf_and_png(pptx_path, out_dir)
    if pdf_path != pptx_path:
        print(f"✅ PDF:  {pdf_path}")
    if png_path:
        print(f"✅ PNG:  {png_path}")


if __name__ == "__main__":
    main()
